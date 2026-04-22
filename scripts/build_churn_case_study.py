from __future__ import annotations

import json
import warnings
from pathlib import Path

import joblib
import sys
import matplotlib

# Only use the headless 'Agg' backend if running as a standard script
if "ipykernel" not in sys.modules:
    matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sklearn.compose import ColumnTransformer
from sklearn.impute import SimpleImputer
from sklearn.inspection import permutation_importance
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import (
    average_precision_score,
    brier_score_loss,
    classification_report,
    precision_recall_curve,
    roc_auc_score,
    roc_curve,
)
from sklearn.pipeline import Pipeline
from sklearn.preprocessing import OneHotEncoder, RobustScaler
from sklearn.tree import DecisionTreeClassifier

warnings.filterwarnings("ignore", category=FutureWarning)

PROJECT_DIR = Path(__file__).resolve().parents[1] if "__file__" in globals() else Path("/Users/ejazanwar/Desktop/Codex/Expedia_Project")
DATA_PATH = Path("/Users/ejazanwar/Downloads/Case Study/PIP_case_study_data 3 (1).csv")
OUT_DIR = PROJECT_DIR / "outputs"
CHART_DIR = OUT_DIR / "charts"
TABLE_DIR = OUT_DIR / "tables"
DECK_DIR = OUT_DIR / "deck"
MODEL_DIR = OUT_DIR / "model"

TARGET = "churn_flag"
ID_COLUMNS = ["email_address", "booking_id"]
LEAKAGE_COLUMNS = ["cancel_date"]
DATE_COLUMNS = ["bk_date", "cancel_date"]

NUMERIC_FEATURES = [
    "coupon_flag",
    "pay_now_flag",
    "cancel_flag",
    "loyalty_tier",
    "total_visit_minutes",
    "total_visit_pages",
    "landing_pages_count",
    "search_pages_count",
    "property_pages_count",
    "bkg_confirmation_pages_count",
    "bounce_visits_count",
    "searched_destinations_count",
    "hotel_star_rating",
    "booking_month",
    "booking_dayofweek",
]

CATEGORICAL_FEATURES = [
    "customer_type",
    "platform",
    "marketing_channel",
]

FEATURES = NUMERIC_FEATURES + CATEGORICAL_FEATURES

COLORS = {
    "ink": "233142",
    "blue": "2F80ED",
    "green": "219653",
    "red": "D64545",
    "gold": "C99700",
    "gray": "6B7280",
    "light": "F3F6FA",
}


def ensure_dirs() -> None:
    for directory in [CHART_DIR, TABLE_DIR, DECK_DIR, MODEL_DIR]:
        directory.mkdir(parents=True, exist_ok=True)


def pct(value: float) -> str:
    return f"{value:.1%}"


def savefig(path: Path) -> None:
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()


def load_and_clean() -> pd.DataFrame:
    df = pd.read_csv(DATA_PATH)
    df.attrs["source_column_count"] = int(df.shape[1])
    df["bk_date"] = pd.to_datetime(df["bk_date"], errors="coerce")
    df["cancel_date"] = pd.to_datetime(df["cancel_date"], errors="coerce")
    df["cancel_flag"] = df["cancel_flag"].fillna(0).astype(int)
    df["marketing_channel"] = df["marketing_channel"].fillna("Missing")
    df["total_visit_minutes"] = df["total_visit_minutes"].fillna(df["total_visit_minutes"].median())
    df["booking_month"] = df["bk_date"].dt.month.astype(int)
    df["booking_dayofweek"] = df["bk_date"].dt.dayofweek.astype(int)
    return df


def validate_data(df: pd.DataFrame) -> dict:
    validation = {
        "row_count": int(len(df)),
        "source_column_count": int(df.attrs.get("source_column_count", df.shape[1])),
        "modeling_column_count_after_engineering": int(df.shape[1]),
        "customer_count": int(df["email_address"].nunique()),
        "duplicate_booking_ids": int(df["booking_id"].duplicated().sum()),
        "date_min": str(df["bk_date"].min().date()),
        "date_max": str(df["bk_date"].max().date()),
        "target_rate": float(df[TARGET].mean()),
        "missing_values_after_cleaning": {
            col: int(df[col].isna().sum())
            for col in [
                "cancel_flag",
                "marketing_channel",
                "total_visit_minutes",
                "bk_date",
                TARGET,
            ]
        },
        "identifier_columns_in_model": sorted(set(ID_COLUMNS).intersection(FEATURES)),
        "leakage_columns_in_model": sorted(set(LEAKAGE_COLUMNS).intersection(FEATURES)),
    }
    assert validation["duplicate_booking_ids"] == 0
    assert not validation["identifier_columns_in_model"]
    assert not validation["leakage_columns_in_model"]
    assert df[TARGET].isin([0, 1]).all()
    (TABLE_DIR / "data_validation.json").write_text(json.dumps(validation, indent=2), encoding="utf-8")
    return validation


def churn_table(df: pd.DataFrame, column: str, labels: dict | None = None) -> pd.DataFrame:
    table = (
        df.groupby(column, dropna=False)[TARGET]
        .agg(churn_rate="mean", bookings="count")
        .reset_index()
        .sort_values("churn_rate", ascending=False)
    )
    table[column] = table[column].map(labels).fillna(table[column]) if labels else table[column]
    table["churn_rate_pct"] = table["churn_rate"].map(lambda x: round(x * 100, 1))
    table.to_csv(TABLE_DIR / f"churn_by_{column}.csv", index=False)
    return table


def plot_churn_bar(table: pd.DataFrame, column: str, title: str, filename: str) -> Path:
    path = CHART_DIR / filename
    plot_df = table.copy()
    plot_df[column] = plot_df[column].astype(str)
    plt.figure(figsize=(9, 4.8))
    sns.barplot(data=plot_df, x="churn_rate", y=column, color=f"#{COLORS['blue']}")
    plt.axvline(plot_df["churn_rate"].mean(), color=f"#{COLORS['gray']}", linestyle="--", linewidth=1)
    plt.title(title, fontsize=14, weight="bold")
    plt.xlabel("Churn rate")
    plt.ylabel("")
    plt.gca().xaxis.set_major_formatter(lambda x, _: f"{x:.0%}")
    for index, row in plot_df.reset_index(drop=True).iterrows():
        plt.text(row["churn_rate"] + 0.01, index, f"{row['churn_rate']:.1%}", va="center", fontsize=9)
    savefig(path)
    return path


def plot_metric_bar(metrics: pd.DataFrame) -> Path:
    path = CHART_DIR / "model_metrics.png"
    plot_df = metrics.melt(id_vars="model", value_vars=["roc_auc", "pr_auc", "brier_score"])
    plt.figure(figsize=(8.5, 4.8))
    sns.barplot(data=plot_df, x="variable", y="value", hue="model", palette=["#2F80ED", "#219653"])
    plt.title("Model validation performance", fontsize=14, weight="bold")
    plt.xlabel("")
    plt.ylabel("Score")
    plt.ylim(0, max(1.0, plot_df["value"].max() * 1.15))
    plt.legend(title="")
    savefig(path)
    return path


def plot_lift(deciles: pd.DataFrame) -> Path:
    path = CHART_DIR / "lift_by_decile.png"
    plt.figure(figsize=(9, 4.8))
    sns.barplot(data=deciles, x="risk_decile", y="lift", color=f"#{COLORS['green']}")
    plt.axhline(1.0, color=f"#{COLORS['gray']}", linestyle="--", linewidth=1)
    plt.title("Lift by predicted-risk decile", fontsize=14, weight="bold")
    plt.xlabel("Risk decile, 10 = highest predicted churn risk")
    plt.ylabel("Lift vs. average churn")
    savefig(path)
    return path


def plot_roc(y_true: pd.Series, y_score: np.ndarray) -> Path:
    path = CHART_DIR / "roc_curve.png"
    fpr, tpr, _ = roc_curve(y_true, y_score)
    auc = roc_auc_score(y_true, y_score)
    plt.figure(figsize=(6, 5))
    plt.plot(fpr, tpr, color=f"#{COLORS['blue']}", linewidth=2, label=f"ROC-AUC {auc:.3f}")
    plt.plot([0, 1], [0, 1], color=f"#{COLORS['gray']}", linestyle="--")
    plt.title("Holdout ROC curve", fontsize=14, weight="bold")
    plt.xlabel("False positive rate")
    plt.ylabel("True positive rate")
    plt.legend(loc="lower right")
    savefig(path)
    return path


def plot_pr(y_true: pd.Series, y_score: np.ndarray) -> Path:
    path = CHART_DIR / "precision_recall_curve.png"
    precision, recall, _ = precision_recall_curve(y_true, y_score)
    ap = average_precision_score(y_true, y_score)
    plt.figure(figsize=(6, 5))
    plt.plot(recall, precision, color=f"#{COLORS['green']}", linewidth=2, label=f"PR-AUC {ap:.3f}")
    plt.title("Holdout precision-recall curve", fontsize=14, weight="bold")
    plt.xlabel("Recall")
    plt.ylabel("Precision")
    plt.legend(loc="upper right")
    savefig(path)
    return path


def plot_importance(importance: pd.DataFrame) -> Path:
    path = CHART_DIR / "driver_importance.png"
    plot_df = importance.head(12).sort_values("importance_mean")
    plt.figure(figsize=(9, 5.4))
    sns.barplot(data=plot_df, x="importance_mean", y="feature", color=f"#{COLORS['blue']}")
    plt.title("Most important model drivers", fontsize=14, weight="bold")
    plt.xlabel("Permutation importance, ROC-AUC decrease")
    plt.ylabel("")
    savefig(path)
    return path


def plot_shap(shap_values: pd.DataFrame) -> Path:
    path = CHART_DIR / "shap_linear_importance.png"
    plot_df = shap_values.head(12).sort_values("mean_abs_shap")
    plt.figure(figsize=(9, 5.4))
    sns.barplot(data=plot_df, x="mean_abs_shap", y="encoded_feature", color=f"#{COLORS['green']}")
    plt.title("SHAP support view for logistic model", fontsize=14, weight="bold")
    plt.xlabel("Mean absolute SHAP value")
    plt.ylabel("")
    savefig(path)
    return path


def plot_risk_bands(risk_bands: pd.DataFrame) -> Path:
    path = CHART_DIR / "risk_bands.png"
    order = ["Low", "Medium", "High"]
    plot_df = risk_bands.set_index("risk_band").loc[order].reset_index()
    plt.figure(figsize=(8.5, 4.8))
    sns.barplot(data=plot_df, x="risk_band", y="actual_churn_rate", color=f"#{COLORS['red']}")
    plt.title("Observed churn by model risk band", fontsize=14, weight="bold")
    plt.xlabel("")
    plt.ylabel("Observed churn rate")
    plt.gca().yaxis.set_major_formatter(lambda x, _: f"{x:.0%}")
    for index, row in plot_df.iterrows():
        plt.text(index, row["actual_churn_rate"] + 0.01, f"{row['actual_churn_rate']:.1%}", ha="center")
    savefig(path)
    return path


def make_preprocessor() -> ColumnTransformer:
    numeric_pipeline = Pipeline(
        steps=[
            ("imputer", SimpleImputer(strategy="median")),
            ("scaler", RobustScaler()),
        ]
    )
    categorical_pipeline = Pipeline(
        steps=[
            ("imputer", SimpleImputer(strategy="most_frequent")),
            ("onehot", OneHotEncoder(handle_unknown="ignore", sparse_output=True)),
        ]
    )
    return ColumnTransformer(
        transformers=[
            ("numeric", numeric_pipeline, NUMERIC_FEATURES),
            ("categorical", categorical_pipeline, CATEGORICAL_FEATURES),
        ]
    )


def time_split(df: pd.DataFrame) -> tuple[pd.DataFrame, pd.DataFrame, pd.Timestamp]:
    ordered = df.sort_values("bk_date").reset_index(drop=True)
    split_index = int(len(ordered) * 0.8)
    cutoff = ordered.loc[split_index, "bk_date"]
    train = ordered[ordered["bk_date"] < cutoff].copy()
    test = ordered[ordered["bk_date"] >= cutoff].copy()
    return train, test, cutoff


def evaluate_model(name: str, model: Pipeline, x_test: pd.DataFrame, y_test: pd.Series) -> dict:
    scores = model.predict_proba(x_test)[:, 1]
    return {
        "model": name,
        "roc_auc": roc_auc_score(y_test, scores),
        "pr_auc": average_precision_score(y_test, scores),
        "brier_score": brier_score_loss(y_test, scores),
        "mean_predicted_churn": float(scores.mean()),
        "actual_churn": float(y_test.mean()),
    }


def build_models(train: pd.DataFrame, test: pd.DataFrame) -> tuple[Pipeline, pd.DataFrame, np.ndarray]:
    x_train, y_train = train[FEATURES], train[TARGET]
    x_test, y_test = test[FEATURES], test[TARGET]

    logistic = Pipeline(
        steps=[
            ("preprocessor", make_preprocessor()),
            ("model", LogisticRegression(max_iter=1000, class_weight="balanced", solver="lbfgs")),
        ]
    )
    tree = Pipeline(
        steps=[
            ("preprocessor", make_preprocessor()),
            ("model", DecisionTreeClassifier(max_depth=6, min_samples_leaf=1000, random_state=42, class_weight="balanced")),
        ]
    )

    logistic.fit(x_train, y_train)
    tree.fit(x_train, y_train)

    metrics = pd.DataFrame(
        [
            evaluate_model("Explainable logistic regression", logistic, x_test, y_test),
            evaluate_model("Shallow decision tree challenger", tree, x_test, y_test),
        ]
    )
    metrics.to_csv(TABLE_DIR / "model_metrics.csv", index=False)
    scores = logistic.predict_proba(x_test)[:, 1]

    report = classification_report(y_test, (scores >= 0.5).astype(int), output_dict=True)
    (TABLE_DIR / "classification_report.json").write_text(json.dumps(report, indent=2), encoding="utf-8")
    joblib.dump(logistic, MODEL_DIR / "churn_logistic_pipeline.joblib")
    return logistic, metrics, scores


def coefficient_table(model: Pipeline) -> pd.DataFrame:
    preprocessor = model.named_steps["preprocessor"]
    feature_names = preprocessor.get_feature_names_out()
    coefficients = model.named_steps["model"].coef_[0]
    coef_df = pd.DataFrame({"encoded_feature": feature_names, "coefficient": coefficients})
    coef_df["odds_ratio"] = np.exp(coef_df["coefficient"])
    coef_df["absolute_coefficient"] = coef_df["coefficient"].abs()
    coef_df = coef_df.sort_values("absolute_coefficient", ascending=False)
    coef_df.to_csv(TABLE_DIR / "logistic_coefficients.csv", index=False)
    return coef_df


def permutation_table(model: Pipeline, test: pd.DataFrame) -> pd.DataFrame:
    sample = test.sample(n=min(40000, len(test)), random_state=42)
    result = permutation_importance(
        model,
        sample[FEATURES],
        sample[TARGET],
        scoring="roc_auc",
        n_repeats=3,
        random_state=42,
        n_jobs=-1,
    )
    importance = pd.DataFrame(
        {
            "feature": FEATURES,
            "importance_mean": result.importances_mean,
            "importance_std": result.importances_std,
        }
    ).sort_values("importance_mean", ascending=False)
    importance.to_csv(TABLE_DIR / "permutation_importance.csv", index=False)
    return importance


def optional_shap_outputs(model: Pipeline, test: pd.DataFrame) -> dict:
    status_path = TABLE_DIR / "shap_status.json"
    try:
        import shap

        sample = test[FEATURES].sample(n=min(2500, len(test)), random_state=42)
        preprocessor = model.named_steps["preprocessor"]
        estimator = model.named_steps["model"]
        transformed = preprocessor.transform(sample)
        if hasattr(transformed, "toarray"):
            transformed = transformed.toarray()
        explainer = shap.LinearExplainer(estimator, transformed)
        values = explainer.shap_values(transformed)
        if isinstance(values, list):
            values = values[-1]
        shap_df = pd.DataFrame(
            {
                "encoded_feature": preprocessor.get_feature_names_out(),
                "mean_abs_shap": np.abs(values).mean(axis=0),
            }
        ).sort_values("mean_abs_shap", ascending=False)
        shap_df.to_csv(TABLE_DIR / "shap_linear_importance.csv", index=False)
        chart = plot_shap(shap_df)
        status = {"status": "created", "rows_sampled": int(len(sample)), "chart": str(chart)}
    except Exception as exc:
        status = {"status": "skipped", "reason": f"{type(exc).__name__}: {exc}"}
    status_path.write_text(json.dumps(status, indent=2), encoding="utf-8")
    return status


def risk_band_outputs(test: pd.DataFrame, scores: np.ndarray) -> tuple[pd.DataFrame, pd.DataFrame]:
    scored = test[["email_address", "bk_date", TARGET, "customer_type", "loyalty_tier", "platform", "marketing_channel"]].copy()
    scored["predicted_churn_probability"] = scores
    scored["risk_band"] = pd.qcut(scores, q=[0, 0.6, 0.85, 1.0], labels=["Low", "Medium", "High"], duplicates="drop")
    scored.to_csv(TABLE_DIR / "holdout_scored_bookings.csv", index=False)

    risk_bands = (
        scored.groupby("risk_band", observed=False)
        .agg(
            bookings=(TARGET, "count"),
            actual_churn_rate=(TARGET, "mean"),
            mean_predicted_churn=("predicted_churn_probability", "mean"),
        )
        .reset_index()
    )
    risk_bands.to_csv(TABLE_DIR / "risk_bands.csv", index=False)

    segments = (
        scored.groupby(["risk_band", "customer_type", "loyalty_tier", "platform"], observed=False)
        .agg(bookings=(TARGET, "count"), actual_churn_rate=(TARGET, "mean"))
        .reset_index()
    )
    segments = segments[segments["bookings"] >= 100].sort_values(["risk_band", "actual_churn_rate", "bookings"], ascending=[False, False, False])
    segments.to_csv(TABLE_DIR / "risk_segments.csv", index=False)
    return risk_bands, segments


def decile_outputs(test: pd.DataFrame, scores: np.ndarray) -> pd.DataFrame:
    deciles = pd.DataFrame({"score": scores, TARGET: test[TARGET].values})
    deciles["risk_decile"] = pd.qcut(deciles["score"], 10, labels=False, duplicates="drop") + 1
    summary = (
        deciles.groupby("risk_decile")
        .agg(bookings=(TARGET, "count"), churn_rate=(TARGET, "mean"), mean_score=("score", "mean"))
        .reset_index()
    )
    summary["lift"] = summary["churn_rate"] / deciles[TARGET].mean()
    summary.to_csv(TABLE_DIR / "lift_by_decile.csv", index=False)
    return summary


def make_eda(df: pd.DataFrame) -> dict[str, Path]:
    labels_loyalty = {0: "Not member", 1: "Base member", 2: "Silver/Gold"}
    labels_star = {0: "Lower star", 1: "Mid star", 2: "High star"}
    outputs = {}
    specs = [
        ("customer_type", None, "New customers churn materially more often", "churn_by_customer_type.png"),
        ("loyalty_tier", labels_loyalty, "Loyalty tier is strongly associated with retention", "churn_by_loyalty_tier.png"),
        ("platform", None, "App bookings show the lowest churn rate", "churn_by_platform.png"),
        ("marketing_channel", None, "Acquisition channel quality varies sharply", "churn_by_marketing_channel.png"),
        ("coupon_flag", {0: "No coupon", 1: "Coupon"}, "Coupon users have lower observed churn", "churn_by_coupon_flag.png"),
        ("pay_now_flag", {0: "Pay later", 1: "Pay now"}, "Pay-now customers are more likely to return", "churn_by_pay_now_flag.png"),
        ("cancel_flag", {0: "Not canceled", 1: "Canceled"}, "Canceled bookings carry slightly higher churn", "churn_by_cancel_flag.png"),
        ("hotel_star_rating", labels_star, "Hotel star rating has a small but visible relationship", "churn_by_hotel_star_rating.png"),
    ]
    for column, labels, title, filename in specs:
        table = churn_table(df, column, labels)
        outputs[column] = plot_churn_bar(table, column, title, filename)
    return outputs


def text_box(slide, text: str, left: float, top: float, width: float, height: float, font_size: int = 18, bold: bool = False, color: str = "233142"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor.from_string(color)
    return box


def bullet_box(slide, bullets: list[str], left: float, top: float, width: float, height: float, font_size: int = 18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for i, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = RGBColor.from_string(COLORS["ink"])
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    text_box(slide, title, 0.55, 0.25, 12.2, 0.45, font_size=24, bold=True)
    if subtitle:
        text_box(slide, subtitle, 0.58, 0.78, 11.8, 0.3, font_size=12, color=COLORS["gray"])
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.08), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor.from_string(COLORS["blue"])
    line.line.color.rgb = RGBColor.from_string(COLORS["blue"])


def add_image(slide, path: Path, left: float, top: float, width: float, height: float | None = None):
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height) if height else None)


def add_metric_card(slide, label: str, value: str, left: float, top: float, color: str = "2F80ED") -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(2.65), Inches(1.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(COLORS["light"])
    shape.line.color.rgb = RGBColor.from_string(color)
    text_box(slide, value, left + 0.12, top + 0.12, 2.4, 0.38, font_size=22, bold=True, color=color)
    text_box(slide, label, left + 0.12, top + 0.58, 2.4, 0.3, font_size=10, color=COLORS["gray"])


def build_deck(
    validation: dict,
    eda_paths: dict[str, Path],
    metrics: pd.DataFrame,
    importance: pd.DataFrame,
    risk_bands: pd.DataFrame,
    segments: pd.DataFrame,
    chart_paths: dict[str, Path],
    cutoff: pd.Timestamp,
) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    base_churn = validation["target_rate"]
    best = metrics.iloc[0]
    high_band = risk_bands[risk_bands["risk_band"].astype(str) == "High"].iloc[0]
    low_band = risk_bands[risk_bands["risk_band"].astype(str) == "Low"].iloc[0]
    top_drivers = importance.head(5)["feature"].tolist()
    top_segment = segments.head(1).iloc[0]

    slide = prs.slides.add_slide(blank)
    text_box(slide, "Expedia churn case study", 0.7, 0.75, 11.8, 0.7, 34, True, COLORS["ink"])
    text_box(slide, "Explainable booking-level model and retention actions", 0.75, 1.55, 11.6, 0.5, 18, False, COLORS["gray"])
    add_metric_card(slide, "Historic booking-level churn", pct(base_churn), 0.8, 2.65, COLORS["red"])
    add_metric_card(slide, "Holdout ROC-AUC", f"{best['roc_auc']:.3f}", 3.75, 2.65, COLORS["blue"])
    add_metric_card(slide, "High-risk band churn", pct(high_band["actual_churn_rate"]), 6.7, 2.65, COLORS["red"])
    bullet_box(
        slide,
        [
            "Churn is concentrated among new, non-loyalty, mobile-web and paid/acquired-channel bookings.",
            "Loyalty membership, app engagement, direct traffic, coupons, and pay-now behavior are associated with stronger retention.",
            "A score-based operating model can focus interventions on the highest-risk bookings after confirmation.",
        ],
        0.85,
        4.05,
        11.8,
        1.6,
        18,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Business question and data", "Booking-level historic customer records with a six-month churn flag")
    add_metric_card(slide, "Bookings", f"{validation['row_count']:,}", 0.7, 1.45, COLORS["blue"])
    add_metric_card(slide, "Customers", f"{validation['customer_count']:,}", 3.55, 1.45, COLORS["green"])
    add_metric_card(slide, "Date range", f"{validation['date_min']} to {validation['date_max']}", 6.4, 1.45, COLORS["gold"])
    bullet_box(
        slide,
        [
            "Target: churn_flag = no repeat booking within six months after checkout.",
            "Model excludes customer and booking identifiers from predictors.",
            "cancel_date is excluded from prediction to avoid post-booking leakage.",
            f"Final validation uses a time-based split; holdout starts on {cutoff.date()}.",
        ],
        0.75,
        3.05,
        11.9,
        2.3,
        18,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Churn baseline and strongest raw patterns")
    add_image(slide, eda_paths["customer_type"], 0.55, 1.35, 5.9, 3.0)
    add_image(slide, eda_paths["loyalty_tier"], 6.85, 1.35, 5.9, 3.0)
    bullet_box(
        slide,
        [
            "New customers have the largest raw churn gap.",
            "Silver/Gold members show a materially lower churn rate than base or non-members.",
            "These patterns are large enough to anchor the retention strategy even before modeling.",
        ],
        1.05,
        4.92,
        11.4,
        1.1,
        16,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Most important modeled drivers")
    add_image(slide, chart_paths["importance"], 0.65, 1.25, 7.1, 4.45)
    bullet_box(
        slide,
        [
            f"Top modeled drivers: {', '.join(top_drivers)}.",
            "Permutation importance measures how much validation ROC-AUC drops when each field is shuffled.",
            "Driver ranking aligns with the raw churn patterns, supporting executive interpretation.",
        ],
        8.05,
        1.65,
        4.6,
        3.2,
        17,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Model approach and validation")
    add_image(slide, chart_paths["metrics"], 0.65, 1.35, 6.1, 3.5)
    bullet_box(
        slide,
        [
            "Primary model: regularized logistic regression with one-hot encoded categorical fields.",
            "Challenger: shallow decision tree to check whether simple nonlinear splits improve validation.",
            "The logistic model remains the recommended production starting point because it is easier to explain and operationalize.",
        ],
        7.15,
        1.45,
        5.2,
        3.1,
        17,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Model lift supports targeted action")
    add_image(slide, chart_paths["lift"], 0.65, 1.25, 6.2, 3.6)
    add_image(slide, chart_paths["risk_bands"], 7.0, 1.25, 5.7, 3.6)
    bullet_box(
        slide,
        [
            f"High-risk bookings churn at {pct(high_band['actual_churn_rate'])}, versus {pct(low_band['actual_churn_rate'])} in the low-risk band.",
            "This separation is enough to prioritize retention treatments rather than treating all bookings equally.",
        ],
        1.0,
        5.25,
        11.4,
        0.9,
        16,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Customer segments most at risk")
    segment_text = (
        f"Highest-risk validation segment: {top_segment['customer_type']} customers, "
        f"loyalty tier {int(top_segment['loyalty_tier'])}, {top_segment['platform']} platform, "
        f"with {int(top_segment['bookings']):,} bookings and {pct(top_segment['actual_churn_rate'])} churn."
    )
    add_image(slide, eda_paths["platform"], 0.65, 1.25, 5.9, 3.25)
    add_image(slide, eda_paths["marketing_channel"], 6.9, 1.25, 5.9, 3.25)
    bullet_box(slide, [segment_text, "Risk segmentation should be refreshed as acquisition mix and customer behavior change."], 0.95, 5.0, 11.6, 1.0, 16)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Recommended actions")
    bullet_box(
        slide,
        [
            "Increase loyalty enrollment prompts for new and non-member customers immediately after booking.",
            "Use risk scores to trigger targeted save offers, app onboarding, or service reassurance for high-risk bookings.",
            "Audit paid and meta acquisition channels for retention quality, not only booking volume.",
            "Promote direct and app re-engagement journeys for customers acquired through high-churn channels.",
            "Test coupon and pay-now nudges with incrementality measurement before scaling.",
        ],
        0.95,
        1.45,
        11.65,
        4.4,
        20,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "How to operationalize churn scoring")
    bullet_box(
        slide,
        [
            "Score each booking after confirmation using available booking, customer, channel, platform, and session attributes.",
            "Route high-risk bookings into retention treatments with clear treatment eligibility and contact-frequency limits.",
            "Track outcomes by scored risk band, treatment group, channel, platform, and loyalty tier.",
            "Recalibrate or retrain monthly until score distributions and intervention lift stabilize.",
        ],
        0.95,
        1.45,
        11.65,
        4.1,
        20,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Caveats and next steps")
    bullet_box(
        slide,
        [
            "The dataset is simulated, so recommendations should be validated through live experiments before rollout.",
            "Booking-level churn can contain repeated customers; production scoring should include customer-history features where available.",
            "cancel_flag is useful after cancellation is known; use a pre-cancellation model if decisions must happen earlier.",
            "Next step: run a holdout campaign or A/B test against the high-risk segment and measure incremental repeat bookings.",
        ],
        0.95,
        1.45,
        11.65,
        4.1,
        20,
    )

    path = DECK_DIR / "Expedia_Churn_Case_Study_Executive_Deck.pptx"
    prs.save(path)
    return path


def write_readme(validation: dict, metrics: pd.DataFrame, deck_path: Path) -> None:
    best = metrics.iloc[0]
    content = f"""# Expedia Churn Case Study

This folder contains the full reproducible case-study package.

## Main Deliverables

- Executive PowerPoint: `{deck_path.relative_to(PROJECT_DIR)}`
- Reproducible analysis script: `scripts/build_churn_case_study.py`
- Charts: `outputs/charts/`
- Tables and validation outputs: `outputs/tables/`
- Saved model pipeline: `outputs/model/churn_logistic_pipeline.joblib`

## Reproduce

```bash
cd /Users/ejazanwar/Desktop/Codex/Expedia_Project
python3 -m venv .venv
.venv/bin/pip install -r requirements.txt
.venv/bin/python scripts/build_churn_case_study.py
```

## Data Summary

- Rows: {validation['row_count']:,}
- Customers: {validation['customer_count']:,}
- Source columns: {validation['source_column_count']:,}
- Booking dates: {validation['date_min']} to {validation['date_max']}
- Historic churn rate: {pct(validation['target_rate'])}

## Primary Model

- Model: regularized logistic regression
- Holdout ROC-AUC: {best['roc_auc']:.3f}
- Holdout PR-AUC: {best['pr_auc']:.3f}
- Brier score: {best['brier_score']:.3f}

## Leakage Controls

- Excluded identifiers: `email_address`, `booking_id`
- Excluded post-booking date field: `cancel_date`
- Final validation uses a time-based split by `bk_date`
"""
    (PROJECT_DIR / "README.md").write_text(content, encoding="utf-8")


def main() -> None:
    ensure_dirs()
    sns.set_theme(style="whitegrid")

    df = load_and_clean()
    validation = validate_data(df)
    eda_paths = make_eda(df)

    train, test, cutoff = time_split(df)
    split_summary = {
        "cutoff_date": str(cutoff.date()),
        "train_rows": int(len(train)),
        "test_rows": int(len(test)),
        "train_churn_rate": float(train[TARGET].mean()),
        "test_churn_rate": float(test[TARGET].mean()),
    }
    (TABLE_DIR / "time_split.json").write_text(json.dumps(split_summary, indent=2), encoding="utf-8")

    model, metrics, scores = build_models(train, test)
    coef_df = coefficient_table(model)
    importance = permutation_table(model, test)
    shap_status = optional_shap_outputs(model, test)
    risk_bands, segments = risk_band_outputs(test, scores)
    deciles = decile_outputs(test, scores)

    chart_paths = {
        "metrics": plot_metric_bar(metrics),
        "lift": plot_lift(deciles),
        "roc": plot_roc(test[TARGET], scores),
        "pr": plot_pr(test[TARGET], scores),
        "importance": plot_importance(importance),
        "risk_bands": plot_risk_bands(risk_bands),
    }

    deck_path = build_deck(validation, eda_paths, metrics, importance, risk_bands, segments, chart_paths, cutoff)
    write_readme(validation, metrics, deck_path)

    summary = {
        "deck": str(deck_path),
        "readme": str(PROJECT_DIR / "README.md"),
        "top_permutation_drivers": importance.head(10).to_dict(orient="records"),
        "top_logistic_coefficients": coef_df.head(10).to_dict(orient="records"),
        "shap_status": shap_status,
        "metrics": metrics.to_dict(orient="records"),
        "split": split_summary,
    }
    (OUT_DIR / "case_study_summary.json").write_text(json.dumps(summary, indent=2, default=str), encoding="utf-8")
    print(json.dumps(summary, indent=2, default=str))


if __name__ == "__main__":
    main()
